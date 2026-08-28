"""화면설계서.pptx 생성기.

실제 브라우저에서 캡처한 screens/*.png 위에 번호 배지가 이미 찍혀 있고,
이 스크립트는 그 번호에 대응하는 디스크립션을 붙여 슬라이드를 만듭니다.
본문 데이터는 screen_spec_data.py 가 갖습니다.

    python docs/planning/build_screen_spec.py
"""

from __future__ import annotations

import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from screen_spec_data import SCREENS

BASE = Path(__file__).resolve().parent
SHOTS = BASE / "screens"
OUT = BASE / "화면설계서.pptx"

PRODUCT = "SalesLuv"
DOCNAME = "화면설계서"
WRITTEN = "2026-08-28"
VERSION = "v1.0"

BRAND = RGBColor(0x1B, 0x6E, 0xF3)
BRAND_DK = RGBColor(0x10, 0x2A, 0x5C)
INK = RGBColor(0x14, 0x1B, 0x2B)
BODY = RGBColor(0x4B, 0x54, 0x66)
MUTED = RGBColor(0x8A, 0x93, 0xA5)
LINE = RGBColor(0xE3, 0xE7, 0xEE)
TINT = RGBColor(0xF4, 0xF7, 0xFC)
ACCENT = RGBColor(0xE2, 0x3C, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0xBF, 0xD4, 0xFF)

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)


# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    return s


def text(slide, x, y, w, h, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """blocks: [(문자열, pt, bold, color[, 줄간격, 문단앞여백pt])]"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, blk in enumerate(blocks):
        body, size, bold, color = blk[0], blk[1], blk[2], blk[3]
        spacing = blk[4] if len(blk) > 4 else 1.0
        before = blk[5] if len(blk) > 5 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_before = Pt(before)
        r = p.add_run()
        r.text = body
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return box


def badge(slide, x, y, n, d=Inches(0.24), fill=ACCENT, size=9):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    return s


def pill(slide, x, y, label, fill, fg=WHITE, size=8.5):
    w = Inches(0.085 * len(label) + 0.30)
    s = rect(slide, x, y, w, Inches(0.24), fill=fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    r.font.name = FONT
    return s, w


def footer(slide, left, page):
    rect(slide, Inches(0), Inches(7.14), W, Inches(0.014), fill=LINE)
    text(slide, Inches(0.45), Inches(7.24), Inches(9), Inches(0.22),
         [(left, 8, False, MUTED)])
    text(slide, Inches(10.9), Inches(7.24), Inches(2.0), Inches(0.22),
         [(f"{PRODUCT} {DOCNAME} · {page}", 8, False, MUTED)], align=PP_ALIGN.RIGHT)


def section_head(slide, kicker, title, desc=None):
    rect(slide, Inches(0), Inches(0), W, Inches(0.92), fill=BRAND_DK)
    text(slide, Inches(0.45), Inches(0.19), Inches(11), Inches(0.24),
         [(kicker, 9.5, True, SKY)])
    text(slide, Inches(0.45), Inches(0.44), Inches(11), Inches(0.34),
         [(title, 17, True, WHITE)])
    if desc:
        text(slide, Inches(0.45), Inches(1.10), Inches(12.4), Inches(0.34),
             [(desc, 10, False, BODY)])


def table(slide, x, y, w, cols, rows, widths, head_h=0.32, row_h=0.30, size=8.5):
    tbl = slide.shapes.add_table(len(rows) + 1, len(cols), x, y, w,
                                 Inches(head_h + row_h * len(rows))).table
    total = sum(widths)
    for i, ratio in enumerate(widths):
        tbl.columns[i].width = Emu(int(w * ratio / total))
    tbl.rows[0].height = Inches(head_h)
    for i in range(len(rows)):
        tbl.rows[i + 1].height = Inches(row_h)

    for c, label in enumerate(cols):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_DK
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else TINT
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 0.95
            r = p.add_run()
            r.text = val
            r.font.size = Pt(size)
            r.font.bold = ci == 0
            r.font.color.rgb = INK if ci == 0 else BODY
            r.font.name = FONT
    return tbl


# ── 화면별 상세 슬라이드 ──────────────────────────────────────────────────
IMG_X, IMG_Y = Inches(0.42), Inches(1.02)
IMG_W = Inches(7.85)
COL_X, COL_W = Inches(8.52), Inches(4.40)


def fit(path, box_w, box_h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    return int(iw * scale), int(ih * scale)


def detail_slide(prs, sc, page):
    slide = blank(prs)
    rect(slide, Inches(0), Inches(0), W, Inches(0.78), fill=BRAND_DK)

    x = Inches(0.45)
    _, pw = pill(slide, x, Inches(0.26), sc["id"], BRAND)
    x += pw + Inches(0.14)
    text(slide, x, Inches(0.24), Inches(6.4), Inches(0.3),
         [(sc["name"], 14.5, True, WHITE)])
    if sc.get("tag"):
        pill(slide, Inches(6.45), Inches(0.28), sc["tag"],
             RGBColor(0x2C, 0x44, 0x74), SKY, 8)

    right = []
    if sc.get("path"):
        right.append((sc["path"], 9.5, True, SKY))
    if sc.get("entry"):
        right.append(("진입 · " + sc["entry"], 8.5, False, RGBColor(0x93, 0xA3, 0xC4), 1.0, 2))
    if right:
        text(slide, Inches(7.7), Inches(0.15), Inches(5.18), Inches(0.5), right,
             align=PP_ALIGN.RIGHT)

    states = sc.get("states") or []
    box_h = Inches(4.95) if states else Inches(5.95)
    iw, ih = fit(SHOTS / sc["img"], IMG_W, box_h)
    ix = IMG_X + Emu(int((IMG_W - iw) / 2))
    rect(slide, ix - Emu(9525), IMG_Y - Emu(9525), Emu(iw + 19050), Emu(ih + 19050),
         fill=None, line=LINE, lw=1)
    slide.shapes.add_picture(str(SHOTS / sc["img"]), ix, IMG_Y, Emu(iw), Emu(ih))

    if states:
        rows = sum(max(1, math.ceil((len(l) + len(b) + 4) / 68)) for l, b in states)
        sy = IMG_Y + Emu(ih) + Inches(0.16)
        sh = Inches(0.22 + 0.20 * rows + 0.045 * (len(states) - 1))
        sh = min(sh, Inches(7.00) - sy)
        rect(slide, IMG_X, sy, IMG_W, sh, fill=TINT, line=LINE, lw=0.75)
        text(slide, IMG_X + Inches(0.15), sy + Inches(0.10), Inches(1.4), Inches(0.2),
             [("상태 · 예외", 8.5, True, BRAND_DK)])
        blocks = []
        for i, (label, body) in enumerate(states):
            blocks.append((f"[{label}]  {body}", 8.5, False, BODY, 1.16, 0 if i == 0 else 3))
        text(slide, IMG_X + Inches(1.48), sy + Inches(0.10), IMG_W - Inches(1.64),
             sh - Inches(0.18), blocks)

    # 오른쪽 디스크립션
    y = Inches(1.00)
    limit = Inches(7.02)
    n = len(sc["notes"])
    size = 9.0 if n <= 6 else 8.4
    lead = 0.185 if n <= 6 else 0.170
    per = 30 if n <= 6 else 32
    for num, head, lines in sc["notes"]:
        badge(slide, COL_X, y + Inches(0.015), num, Inches(0.23), size=8.5)
        head_rows = max(1, math.ceil(len(head) / 22))
        text(slide, COL_X + Inches(0.32), y, COL_W - Inches(0.32),
             Inches(0.21 * head_rows), [(head, size + 0.6, True, INK, 1.1)])
        y += Inches(0.225 * head_rows)
        for ln in lines:
            rows = max(1, math.ceil(len(ln) / per))
            text(slide, COL_X + Inches(0.32), y, COL_W - Inches(0.32),
                 Inches(lead * rows), [("· " + ln, size, False, BODY, 1.16)])
            y += Inches(lead * rows)
        y += Inches(0.105)
        if y > limit:
            break

    footer(slide, f"{sc['id']} · {sc['name']}", page)
    return slide


# ── 앞부분 ────────────────────────────────────────────────────────────────
def cover(prs):
    slide = blank(prs)
    rect(slide, Inches(0), Inches(0), W, H, fill=BRAND_DK)
    rect(slide, Inches(0), Inches(0), Inches(0.16), H, fill=BRAND)
    text(slide, Inches(1.1), Inches(2.05), Inches(9), Inches(0.4),
         [("SALES HISTORY WORKSPACE", 11, True, SKY)])
    text(slide, Inches(1.1), Inches(2.50), Inches(10), Inches(1.1),
         [(PRODUCT, 54, True, WHITE)])
    text(slide, Inches(1.1), Inches(3.62), Inches(10), Inches(0.7),
         [(DOCNAME, 30, True, RGBColor(0x9E, 0xC2, 0xFF))])
    rect(slide, Inches(1.13), Inches(4.55), Inches(2.2), Inches(0.03), fill=BRAND)
    text(slide, Inches(1.1), Inches(4.85), Inches(10), Inches(1.0), [
        ("실제 구현 화면을 캡처해 UI 번호와 동작·상태·예외를 붙인 문서입니다.", 12, False, RGBColor(0xC6, 0xD3, 0xE8)),
        ("기준 문서 · docs/planning/유저 시나리오.md", 10, False, RGBColor(0x8A, 0x9C, 0xBC), 1.0, 8),
    ])
    text(slide, Inches(1.1), Inches(6.35), Inches(10), Inches(0.5), [
        (f"작성일 {WRITTEN}    ·    버전 {VERSION}    ·    확인 환경 로컬 (1600×1000)",
         10.5, True, RGBColor(0xB8, 0xCA, 0xE8)),
    ])
    return slide


def how_to_read(prs, page):
    slide = blank(prs)
    section_head(slide, "00", "문서 읽는 법",
                 "화면 1개 = 슬라이드 1개. 왼쪽은 실제 캡처, 오른쪽은 그 번호에 대응하는 동작 설명입니다.")

    # 레이아웃 모형
    ox, oy = Inches(0.45), Inches(1.62)
    rect(slide, ox, oy, Inches(6.2), Inches(0.55), fill=BRAND_DK)
    pill(slide, ox + Inches(0.16), oy + Inches(0.15), "CU-01", BRAND)
    text(slide, ox + Inches(1.0), oy + Inches(0.14), Inches(3), Inches(0.28),
         [("화면명", 11, True, WHITE)])
    text(slide, ox + Inches(4.1), oy + Inches(0.10), Inches(2.0), Inches(0.4),
         [("/경로", 9, True, SKY), ("진입 · 진입 경로", 7.5, False, MUTED, 1.0, 1)],
         align=PP_ALIGN.RIGHT)
    rect(slide, ox, oy + Inches(0.70), Inches(3.75), Inches(2.35), fill=TINT, line=LINE)
    text(slide, ox, oy + Inches(1.75), Inches(3.75), Inches(0.3),
         [("실제 화면 캡처 + 번호 배지", 9.5, True, MUTED)], align=PP_ALIGN.CENTER)
    rect(slide, ox, oy + Inches(3.14), Inches(3.75), Inches(0.62), fill=TINT, line=LINE)
    text(slide, ox + Inches(0.12), oy + Inches(3.22), Inches(3.5), Inches(0.45),
         [("상태 · 예외", 8.5, True, BRAND_DK),
          ("Empty / Loading / Error / Validation / 조건부 노출", 8, False, BODY, 1.0, 2)])
    for i in range(3):
        yy = oy + Inches(0.78 + 0.72 * i)
        badge(slide, ox + Inches(3.95), yy, i + 1, Inches(0.23), size=8.5)
        text(slide, ox + Inches(4.27), yy - Inches(0.02), Inches(1.9), Inches(0.55),
             [("요소 이름", 9.5, True, INK), ("· 동작 · 정책 설명", 8.5, False, BODY, 1.1, 2)])

    # 규칙
    rx = Inches(7.15)
    text(slide, rx, Inches(1.62), Inches(5.7), Inches(0.3), [("작성 규칙", 12, True, INK)])
    rules = [
        ("번호 배지", "캡처 위 빨간 원과 오른쪽 설명의 번호는 1:1로 맞춥니다. 배지가 없는 요소는 설명하지 않습니다."),
        ("설명 기준", "\"무엇인지\"가 아니라 \"어떻게 동작하는지 · 왜 그런지\"를 씁니다. 모든 요소를 설명하지 않습니다."),
        ("상태 표기", "해피 패스만 그리지 않고 Empty / Loading / Error / Validation / 조건부 노출을 따로 적습니다."),
        ("화면 ID", "도메인 약어 + 두 자리. 슬라이드 번호가 아니라 이 ID로 화면을 지칭합니다."),
        ("근거", "모든 문구·동작은 2026-08-28 로컬 실행 화면에서 직접 확인한 것입니다. 코드만 보고 추측하지 않았습니다."),
    ]
    y = Inches(2.02)
    for label, body in rules:
        rect(slide, rx, y + Inches(0.055), Inches(0.05), Inches(0.30), fill=BRAND)
        text(slide, rx + Inches(0.18), y, Inches(5.5), Inches(0.62),
             [(label, 10, True, INK), (body, 9, False, BODY, 1.18, 2)])
        y += Inches(0.78)

    # ID 체계
    text(slide, rx, Inches(6.02), Inches(5.7), Inches(0.3),
         [("화면 ID 체계", 12, True, INK)])
    ids = "AU 인증   DB 대시보드   CU 고객   CP 고객불만   CA 캘린더·일정   RP 업무보고\nDL 영업   SA 매출   DC 자료실   NT 알림   MY 마이페이지   QT·CT·OD 견적·계약·발주   MG 관리   ER 오류"
    text(slide, rx, Inches(6.40), Inches(5.7), Inches(0.6),
         [(ids, 9, False, BODY, 1.35)])
    footer(slide, "문서 읽는 법", page)
    return slide


def flow_slide(prs, page):
    slide = blank(prs)
    section_head(slide, "01", "전체 사용자 Flow")
    text(slide, Inches(0.45), Inches(1.20), Inches(4.6), Inches(1.0),
         [("유저 시나리오의 업무 흐름을 실제 구현된 화면 ID로 옮긴 것입니다.", 9.5, False, BODY, 1.3),
          ("로그인 뒤 대시보드가 모든 갈래의 출발점입니다.", 9.5, False, BODY, 1.3)])

    def node(x, y, w, h, label, sub, fill=WHITE, fg=INK, border=LINE, bold=True):
        rect(slide, x, y, w, h, fill=fill, line=border, lw=1,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
        blocks = [(label, 9.5, bold, fg)]
        if sub:
            blocks.append((sub, 7.5, False, MUTED if fg == INK else SKY, 1.0, 2))
        text(slide, x, y + Inches(0.13 if sub else 0.20), w, h,
             blocks, align=PP_ALIGN.CENTER)

    def arrow(x1, y1, x2, y2, color=RGBColor(0xB6, 0xC2, 0xD6)):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = color
        c.line.width = Pt(1.25)
        return c

    NW, NH = Inches(1.72), Inches(0.62)
    # 진입
    node(Inches(0.45), Inches(3.30), NW, NH, "로그인", "AU-01", BRAND_DK, WHITE, BRAND_DK)
    node(Inches(0.45), Inches(4.15), NW, NH, "계정 요청 · 비밀번호 설정", "AU-02 · AU-03", TINT)
    arrow(Inches(0.45 + 0.86), Inches(4.15), Inches(0.45 + 0.86), Inches(3.92))

    node(Inches(2.62), Inches(3.30), NW, NH, "대시보드", "DB-01", BRAND, WHITE, BRAND)
    arrow(Inches(2.17), Inches(3.61), Inches(2.62), Inches(3.61))

    branches = [
        ("고객 관리", "CU-01 → CU-02 · CU-03", 1.02),
        ("일정 · 미팅", "CA-01 → CA-02 → RP-02", 1.92),
        ("업무보고", "RP-01 → RP-04 → RP-05 → RP-06", 2.82),
        ("영업 관리", "DL-01 · DL-02 → DL-04", 3.72),
        ("매출 분석", "SA-01", 4.62),
        ("자료 관리", "DC-01 → DC-02 → DC-04", 5.52),
        ("CS 관리", "CP-01 → CP-03 → CP-02", 6.42),
    ]
    bx = Inches(5.30)
    bw = Inches(3.15)
    for label, sub, by in branches:
        node(bx, Inches(by), bw, Inches(0.62), label, sub, TINT)
        arrow(Inches(2.62 + 1.72), Inches(3.61), bx, Inches(by + 0.31))

    # 마무리 열
    node(Inches(9.00), Inches(1.02), Inches(3.9), Inches(0.62),
         "고객 상세 · 등록 (모달 / 드로어)", "CU-02 · CU-03 · CU-04 · CU-05", WHITE)
    node(Inches(9.00), Inches(1.92), Inches(3.9), Inches(0.62),
         "AI 브리핑 → 업무보고서 작성", "CA-02 → RP-02 → RP-03", WHITE)
    node(Inches(9.00), Inches(2.82), Inches(3.9), Inches(0.62),
         "일일 → 주간 → 월간 (앞 단계가 재료)", "RP-04 → RP-05 → RP-06", WHITE)
    node(Inches(9.00), Inches(3.72), Inches(3.9), Inches(0.62),
         "견적 · 계약 · 발주로 이어짐", "QT-01 · CT-01 · OD-01", WHITE)
    node(Inches(9.00), Inches(4.62), Inches(3.9), Inches(0.62),
         "기간 × 회사 · 지역 · 상품", "SA-01 · SA-02", WHITE)
    node(Inches(9.00), Inches(5.52), Inches(3.9), Inches(0.62),
         "딜 · 상품 연결 + 버전 관리", "DC-03 · DC-04", WHITE)
    node(Inches(9.00), Inches(6.42), Inches(3.9), Inches(0.62),
         "접수 → 원인파악 → 처리중 → 처리완료", "CP-02 에서 상태 변경", WHITE)
    for _, _, by in branches:
        arrow(Inches(5.30 + 3.15), Inches(by + 0.31), Inches(9.00), Inches(by + 0.31))

    text(slide, Inches(0.45), Inches(5.35), Inches(4.6), Inches(1.5), [
        ("공통", 10, True, INK),
        ("· 로그아웃 상태로 보호 화면 접근 → /login 으로 이동", 8.5, False, BODY, 1.25, 4),
        ("· 팀원이 팀장 전용 주소 직접 입력 → 대시보드로 되돌림", 8.5, False, BODY, 1.25),
        ("· 정의되지 않은 주소 → 404 (ER-01)", 8.5, False, BODY, 1.25),
        ("· 백엔드 응답 실패 → 연결 안내 모달 (ER-02)", 8.5, False, BODY, 1.25),
    ])
    footer(slide, "전체 사용자 Flow", page)
    return slide


SCREEN_LIST = [
    ("AU-01", "로그인", "/login", "서비스 첫 진입", "이메일 로그인, 인증 실패 안내, 로컬 전용 빠른 로그인"),
    ("AU-02", "회원가입 (계정 발급 요청)", "/signup", "로그인 하단 링크", "이메일 남기고 계정 발급 요청"),
    ("AU-03", "비밀번호 설정", "/set-password", "초대 메일 링크", "8자 이상 검증, 링크 만료 처리"),
    ("DB-01", "대시보드", "/", "로그인 직후", "공지·지시사항, KPI 4종, 주간 일정, 일정 추가"),
    ("DB-02", "공지 · 지시사항 상세", "/ (드로어)", "공지 항목 클릭", "공지 원문 확인"),
    ("DB-03", "C/S 대응요청 목록", "/ (드로어)", "C/S KPI 클릭", "미완료·완료 건수, 긴급도·상태"),
    ("DB-04", "계약갱신 예정", "/ (드로어)", "계약갱신 KPI 클릭", "30일 이내 종료 예정 계약"),
    ("DB-05", "대시보드 (팀장)", "/", "팀장 로그인", "관리 메뉴 노출, 팀 전체 범위 전환"),
    ("CU-01", "고객현황", "/customers", "사이드바 고객", "통합 검색, 컬럼 설정, 엑셀 내보내기, 정렬"),
    ("CU-02", "고객 상세", "/customers (드로어)", "고객 행 클릭", "연락처·회사·고객 정보·메모"),
    ("CU-03", "고객 등록", "/customers (모달)", "고객현황 버튼", "회사 자동완성, 필수값 검증, 유입경로"),
    ("CU-04", "명함으로 고객 등록", "/customers (모달)", "고객현황 버튼", "명함 이미지에서 정보 추출"),
    ("CU-05", "엑셀(CSV) 고객 등록", "/customers (모달)", "고객현황 버튼", "CSV 다량 등록"),
    ("CP-01", "고객불만", "/complaints", "사이드바 고객", "상태 탭·건수, 검색, 목록"),
    ("CP-02", "고객불만 상세", "/complaints (드로어)", "불만 행 클릭", "접수 내용, 답변 이력·등록, 상태 변경"),
    ("CP-03", "고객불만 등록", "/complaints (모달)", "고객불만 버튼", "회사 → 딜 종속 선택, 긴급도"),
    ("CA-01", "캘린더", "/calendar", "사이드바 영업", "월간 그리드, 드래그 이동, AI 추천 일정"),
    ("CA-02", "일정 상세", "/ · /calendar (드로어)", "일정 클릭", "연결 딜, AI 브리핑, 보고서 작성 연결"),
    ("CA-03", "일정 등록", "/calendar (모달)", "날짜 칸 + / 일정 추가", "고객사 → 고객 종속 선택, 미등록 회사 즉시 생성"),
    ("RP-01", "업무보고", "/daily", "사이드바 영업", "종류 탭, 주간·월간 달력, 상태·기간 필터"),
    ("RP-02", "업무보고서 작성", "/meetings/new", "일정 상세 · 일정 행", "미팅 정보, 딜 연결, 첨부, AI 보고서 작성"),
    ("RP-03", "업무보고서 AI 결과", "/meetings/new", "AI 작성 완료", "AI 원본 고정 보관, 본문 직접 수정"),
    ("RP-04", "일일업무보고 작성", "/daily/new", "보고서 작성하기", "그날 활동 선택, 필수 항목, 제출 조건"),
    ("RP-05", "주간업무보고 작성", "/daily/new?kind=weekly", "보고서 작성하기", "제출된 일일업무를 재료로 작성"),
    ("RP-06", "월간업무보고 작성", "/daily/new?kind=monthly", "보고서 작성하기", "제출된 주간업무 + 시스템 집계"),
    ("RP-07", "보고서 대상 일정 고르기", "/daily/pick", "보고서 작성하기", "업무보고서를 붙일 일정 1건 선택"),
    ("RP-08", "업무보고서 상세", "/meetings/:id", "업무보고 목록", "제출된 보고서와 근거 확인"),
    ("DL-01", "영업현황 (리스트)", "/deals", "사이드바 영업", "단계 탭, 파이프라인·기간 필터, 정렬"),
    ("DL-02", "영업현황 (보드)", "/deals/board", "영업현황 › 보드", "단계별 컬럼, 드래그로 단계 변경"),
    ("DL-03", "영업 딜 추가", "/deals (모달)", "영업현황 버튼", "고객사·제품·파이프라인"),
    ("DL-04", "딜 상세", "/deals (드로어)", "딜 행·카드 클릭", "딜 정보 + 견적·계약·발주 연결"),
    ("SA-01", "매출분석", "/sales", "사이드바 · 대시보드 매출 목표", "기간 6종 × 회사·지역·상품별, CSV"),
    ("DC-01", "자료실", "/documents", "사이드바 자료실", "분류 탭, 검색, 버전·연결 표시"),
    ("DC-02", "파일 업로드", "/documents (모달)", "자료실 버튼", "분류·연결 대상(상품/딜)·메모"),
    ("DC-04", "자료 상세", "/documents (드로어)", "자료 행 클릭", "자료 정보, 버전 이력, 새 버전"),
    ("NT-01", "알림", "/notifications", "헤더 벨", "전체 / 읽지 않음"),
    ("MY-01", "마이페이지", "/mypage", "사이드바 이름 · 헤더 아바타", "내 정보, 약관·정책, 로그아웃"),
    ("QT-01", "견적 현황", "/quotes", "딜 상세 견적 섹션", "견적 상태 5단계 (시나리오 외)"),
    ("CT-01", "계약 현황", "/contracts", "딜 상세 계약 섹션", "계약 상태 5단계 (시나리오 외)"),
    ("OD-01", "발주 관리", "/orders", "딜 상세 발주 섹션", "발주 상태 6단계, 발주 상세 (시나리오 외)"),
    ("MG-01", "공지관리", "/notices", "사이드바 관리", "공지·지시사항 등록과 노출 관리 (팀장)"),
    ("MG-02", "상품관리", "/products", "사이드바 관리", "상품 등록·단가·유효기간 (팀장)"),
    ("MG-03", "팀 관리", "/team", "사이드바 관리", "역할·재직·월 목표 (팀장, 저장 미구현)"),
    ("ER-01", "404", "정의되지 않은 주소", "잘못된 주소", "셸 안에서 안내 + 대시보드 복귀"),
    ("ER-02", "서버 연결 불가", "전 화면 공통", "백엔드 응답 실패", "안내 모달, 닫기 / 다시 시도"),
]


def list_slides(prs, page):
    cols = ["화면번호", "화면명", "경로", "진입 경로", "주요 기능"]
    widths = [0.9, 2.3, 2.2, 2.3, 5.3]
    made = []
    chunks = [SCREEN_LIST[:16], SCREEN_LIST[16:32], SCREEN_LIST[32:]]
    for i, chunk in enumerate(chunks):
        slide = blank(prs)
        section_head(slide, "02", f"화면 목록 ({i + 1}/{len(chunks)})",
                     "구현된 화면 45개. 상태·변형 화면(Empty·Loading·Error)은 각 화면의 상세 장에 함께 실었습니다."
                     if i == 0 else None)
        table(slide, Inches(0.45), Inches(1.62 if i == 0 else 1.30), Inches(12.44),
              cols, [list(r) for r in chunk], widths, row_h=0.315, size=8.5)
        footer(slide, f"화면 목록 ({i + 1}/{len(chunks)})", page + i)
        made.append(slide)
    return made


LINKS = [
    ("로그인", "로그인 성공", "대시보드", "AU-01 → DB-01", ""),
    ("로그인", "회원가입", "계정 발급 요청", "AU-01 → AU-02", ""),
    ("초대 메일", "링크 열기", "비밀번호 설정 → 로그인", "AU-03 → AU-01", ""),
    ("대시보드", "매출 목표 클릭", "매출분석", "DB-01 → SA-01", "화면을 옮기는 유일한 KPI"),
    ("대시보드", "일정 클릭", "일정 상세 드로어", "DB-01 → CA-02", ""),
    ("대시보드", "C/S 대응요청 클릭", "C/S 목록 드로어", "DB-01 → DB-03", "시나리오는 화면 이동, 실제는 드로어"),
    ("대시보드", "계약갱신 클릭", "계약갱신 드로어", "DB-01 → DB-04", "시나리오는 화면 이동, 실제는 드로어"),
    ("대시보드", "공지 클릭", "공지 상세 드로어", "DB-01 → DB-02", ""),
    ("캘린더", "일정 클릭", "일정 상세 드로어", "CA-01 → CA-02", "대시보드와 같은 드로어"),
    ("캘린더", "날짜 칸 + ", "일정 등록 모달", "CA-01 → CA-03", "빈 칸 클릭만으로는 열리지 않음"),
    ("일정 상세", "업무보고서 작성", "업무보고서 작성", "CA-02 → RP-02", ""),
    ("업무보고", "보고서 작성하기", "일정 고르기 / 일일·주간·월간", "RP-01 → RP-07 · RP-04~06", ""),
    ("업무보고서 작성", "AI 보고서 작성", "AI 결과 확인 → 사용자 수정", "RP-02 → RP-03", ""),
    ("고객현황", "고객 클릭", "고객 상세 드로어", "CU-01 → CU-02", ""),
    ("고객현황", "고객 등록", "고객 등록 모달", "CU-01 → CU-03", "시나리오는 별도 화면, 실제는 모달"),
    ("고객불만", "불만 클릭", "고객불만 상세 드로어", "CP-01 → CP-02", ""),
    ("고객불만", "고객불만 등록", "고객불만 등록 모달", "CP-01 → CP-03", ""),
    ("영업현황", "딜 클릭", "딜 상세 드로어", "DL-01 · DL-02 → DL-04", ""),
    ("딜 상세", "견적 · 계약 · 발주", "각 현황 화면", "DL-04 → QT-01 · CT-01 · OD-01", "사이드바에 없는 진입 경로"),
    ("자료실", "파일 업로드", "업로드 모달 → 목록 반영", "DC-01 → DC-02", ""),
    ("자료실", "자료 클릭", "자료 상세 · 버전 이력", "DC-01 → DC-04", ""),
    ("전 화면", "권한 없는 주소 입력", "대시보드로 되돌림", "→ DB-01", "팀장 전용 화면 접근 시"),
]


def link_slides(prs, page):
    chunks = [LINKS[:11], LINKS[11:]]
    made = []
    for i, chunk in enumerate(chunks):
        slide = blank(prs)
        section_head(slide, "03", f"화면 간 연결 관계 ({i + 1}/{len(chunks)})",
                     "실제 구현 기준입니다. 시나리오와 다른 부분은 비고에 표시했습니다." if i == 0 else None)
        table(slide, Inches(0.45), Inches(1.62 if i == 0 else 1.30), Inches(12.44),
              ["시작 화면", "사용자 행동", "이동 · 결과", "화면 ID", "비고"],
              [list(r) for r in chunk], [1.9, 2.1, 3.1, 2.6, 3.2], row_h=0.33, size=9)
        footer(slide, f"화면 간 연결 관계 ({i + 1}/{len(chunks)})", page + i)
        made.append(slide)
    return made


VALIDATION = [
    ("AU-01", "이메일 · 비밀번호 미입력", "브라우저 기본 필수 안내", "제출 차단"),
    ("AU-01", "자격 불일치", "이메일 또는 비밀번호가 올바르지 않습니다.", "폼 위에 표시, 입력값 유지"),
    ("AU-03", "비밀번호 8자 미만", "8자 이상 입력해 주세요.", "항목 아래 표시, 제출 버튼 비활성"),
    ("AU-03", "비밀번호 확인 불일치", "두 비밀번호가 다릅니다.", "항목 아래 표시, 제출 버튼 비활성"),
    ("AU-03", "토큰 없음 · 만료", "링크가 만료되었거나 올바르지 않습니다.", "폼을 그리지 않고 재발송 요청 안내"),
    ("CU-03", "이름 미입력", "이름을 입력하세요.", "제출 시 검사, 항목 아래 표시"),
    ("CU-03", "전화 미입력", "전화번호를 입력하세요.", "제출 시 검사, 항목 아래 표시"),
    ("CU-03", "기존 회사 선택", "—", "사업자등록번호·주소가 읽기 전용으로 바뀜"),
    ("CP-03", "회사 미선택", "회사를 먼저 선택하세요", "딜 선택 필드 비활성"),
    ("CA-03", "고객사 미선택", "고객사를 먼저 고르세요", "고객 선택 필드 비활성"),
    ("RP-02", "미팅 내용·첨부 모두 없음", "—", "AI 보고서 작성 버튼 비활성"),
    ("RP-04", "필수 항목 · 자료 부족", "제출 전 확인: 자료 1건 이상", "보고서 제출 버튼 비활성"),
    ("RP-04", "같은 기간에 작성중 보고서 존재", "이 기간에 작성중인 보고서가 있어 이어 씁니다.", "새 보고서를 만들지 않고 기존 건을 엽니다"),
    ("CU-05", "파일 미선택", "파일을 먼저 선택하세요", "실행 버튼 라벨에 표시, 비활성"),
    ("전 화면", "백엔드 응답 실패", "서버에 연결할 수 없습니다", "안내 모달, 닫아도 화면·입력 유지"),
]


def validation_slide(prs, page):
    slide = blank(prs)
    section_head(slide, "04", "입력값 검증 · 처리 시나리오",
                 "화면에서 실제로 확인한 문구만 옮겼습니다. 문구가 없는 조건부 동작은 문구 칸을 '—'로 두었습니다.")
    table(slide, Inches(0.45), Inches(1.62), Inches(12.44),
          ["화면", "상황", "문구", "처리"],
          [list(r) for r in VALIDATION], [1.3, 3.4, 4.2, 4.1], row_h=0.335, size=9)
    footer(slide, "입력값 검증", page)
    return slide


GAPS_MISSING = [
    ("21. 자료실 등록", "AI 에이전트가 자료를 읽고 딜 정보를 제안 → 사용자 승인 → 딜에 저장",
     "업로드·연결·버전 관리까지만 구현. AI 검토·승인 화면 없음"),
    ("2. 대시보드", "C/S · 계약갱신 클릭 시 해당 관리 화면으로 이동",
     "화면 이동 대신 드로어 목록만 엽니다 (DB-03 · DB-04)"),
    ("5. 고객 등록", "고객 등록 전용 화면", "목록 위 모달로 구현 (CU-03)"),
    ("18. 영업 현황 등록", "입력 항목에 유입경로 포함",
     "유입경로는 고객 등록(CU-03)으로 옮겼습니다 — 딜이 아니라 고객의 속성으로 정리"),
    ("16. 월간 업무", "월을 걸친 주(8/30~9/4)의 소속 정책", "시나리오에서도 보류. 구현도 미정"),
]

GAPS_EXTRA = [
    ("7. 고객 불만 자세히 보기", "답변 등록을 '추후 기능'으로 표기", "답변 이력·답변 등록·상태 변경 모두 구현됨 (CP-02)"),
    ("—", "시나리오에 없음", "견적 현황(QT-01) · 계약 현황(CT-01) · 발주 관리(OD-01)"),
    ("—", "시나리오에 없음", "공지관리(MG-01) · 상품관리(MG-02) · 팀 관리(MG-03) — 팀장 전용"),
    ("—", "시나리오에 없음", "보기 범위 스위처(내 현황 / 팀 전체), 자료 버전 이력, 컬럼 설정"),
]

GAPS_CHECK = [
    ("SA-01 · DB-01", "대시보드 KPI는 8월 목표 ₩25.0M(83.2%)인데 매출분석은 '목표 미설정'",
     "두 화면의 목표값 출처를 맞출지 결정 필요"),
    ("DC-01", "분류 탭 건수가 실제 분류를 세지 않음 (상품설명서 1건인데 탭은 0, 전체만 1)",
     "집계 기준 확인 필요"),
    ("MG-03", "'변경은 이 화면 안에만 남습니다' — 저장이 서버에 반영되지 않음", "백엔드 연결 필요"),
    ("NT-01 · DC-01", "알림 0건, 자료실 1건(문서 작성용 샘플)", "데모 데이터 보강 여부 결정 필요"),
]


def gap_slide(prs, page):
    slide = blank(prs)
    section_head(slide, "05", "유저 시나리오와 구현의 차이",
                 "실제 화면을 기준으로 정리했습니다. 구현되지 않은 화면은 이 문서에 싣지 않았습니다.")

    y = Inches(1.48)
    for title, color, rows, cols, widths in [
        ("시나리오에 있으나 화면에 없거나 다르게 구현된 것", ACCENT, GAPS_MISSING,
         ["시나리오 항목", "시나리오 내용", "실제 구현"], [2.4, 4.9, 5.1]),
        ("구현이 시나리오보다 앞선 것", RGBColor(0x0F, 0x7B, 0x4F), GAPS_EXTRA,
         ["시나리오 항목", "시나리오 내용", "실제 구현"], [2.4, 4.9, 5.1]),
        ("확인이 필요한 것", RGBColor(0xC2, 0x7A, 0x0A), GAPS_CHECK,
         ["화면", "관찰한 사실", "필요한 결정"], [2.4, 6.2, 3.8]),
    ]:
        rect(slide, Inches(0.45), y + Inches(0.03), Inches(0.05), Inches(0.20), fill=color)
        text(slide, Inches(0.62), y, Inches(12.2), Inches(0.24), [(title, 10.5, True, INK)])
        y += Inches(0.28)
        table(slide, Inches(0.45), y, Inches(12.44), cols,
              [list(r) for r in rows], widths, head_h=0.24, row_h=0.26, size=8.2)
        y += Inches(0.24 + 0.26 * len(rows) + 0.26)

    footer(slide, "시나리오와 구현의 차이", page)
    return slide


# ── 조립 ──────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    cover(prs)
    page = 2
    how_to_read(prs, page); page += 1
    flow_slide(prs, page); page += 1
    made = list_slides(prs, page); page += len(made)
    made = link_slides(prs, page); page += len(made)

    sec = blank(prs)
    section_head(sec, "06", "화면별 상세 설계",
                 "왼쪽은 실제 브라우저 캡처(1600×1000), 오른쪽은 캡처 위 번호에 대응하는 동작 설명입니다.")
    text(sec, Inches(0.45), Inches(2.2), Inches(12.4), Inches(3.5), [
        ("· 화면 1개 = 슬라이드 1개. 상태·변형 화면은 머리말에 Empty / Loading / Error / 권한 등으로 표시했습니다.", 11, False, BODY, 1.6),
        ("· 캡처 위 빨간 번호와 오른쪽 설명의 번호는 1:1로 대응합니다.", 11, False, BODY, 1.6),
        ("· 화면 아래 '상태 · 예외' 칸에는 그 화면에서 실제로 확인한 상태만 적었습니다.", 11, False, BODY, 1.6),
        (f"· 총 {len(SCREENS)}장.", 11, True, INK, 1.6),
    ])
    footer(sec, "화면별 상세 설계", page); page += 1

    for sc in SCREENS:
        detail_slide(prs, sc, page)
        page += 1

    validation_slide(prs, page); page += 1
    gap_slide(prs, page); page += 1

    prs.save(OUT)
    print(f"저장: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)}장)")


if __name__ == "__main__":
    main()
