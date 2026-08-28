"""화면설계서_Final.pptx 에 사이트맵 슬라이드 한 장만 얹습니다.

기존 69장은 손대지 않습니다. 새 슬라이드를 9번째 자리에 끼워 넣고
화면설계서_Final_사이트맵.pptx 로 따로 저장합니다.

    python docs/planning/add_sitemap.py
"""

from __future__ import annotations

import sys
from pathlib import Path

BASE = Path(__file__).resolve().parent
sys.path.insert(0, str(BASE))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 디자인 토큰과 공용 헬퍼는 기존 생성기에서 그대로 가져다 씁니다.
from build_screen_spec import (
    BODY, BRAND, BRAND_DK, INK, LINE, MUTED, SKY, TINT, WHITE,
    footer, rect, section_head, text,
)

SRC = BASE / "화면설계서_Final.pptx"
OUT = BASE / "화면설계서_Final_사이트맵.pptx"
AT = 8  # 0-based. 8번 슬라이드(03 화면 간 연결 관계 2/2) 뒤 = 9번째 자리
PAGE = 9
WIRE = RGBColor(0xB6, 0xC2, 0xD6)

# ── 좌표 ──────────────────────────────────────────────────────────────────
T1_X, T1_W = Inches(0.45), Inches(1.20)
T2_X, T2_W = Inches(1.80), Inches(1.62)
T3_X, T3_W = Inches(3.58), Inches(9.30)
CARD_W, CARD_H, GAP = Inches(1.748), Inches(0.42), Inches(0.14)

Y0 = 1.36
ROW = 0.42
VGAP = 0.06

# ── 사이트맵 데이터 (02 화면 목록에 실린 45개 화면만) ─────────────────────
# (대분류, 부제, [행]) / 행 = (메뉴 라벨, 메뉴 화면ID, [카드 줄])
# 카드 = (라벨, 화면ID, 꼬리표, 앞에 놓을 구분자)
SECTIONS = [
    ("대시보드", None, [
        ("대시보드", "DB-01", [[
            ("공지 · 지시사항 상세", "DB-02", "드로어", None),
            ("C/S 대응요청 목록", "DB-03", "드로어", "·"),
            ("계약갱신 예정", "DB-04", "드로어", "·"),
            ("대시보드 (팀장)", "DB-05", "팀장 전용", "·"),
        ]]),
    ]),
    ("고객", None, [
        ("고객현황", "CU-01", [[
            ("고객 상세", "CU-02", "드로어", None),
            ("고객 등록", "CU-03", "모달", "·"),
            ("명함으로 고객 등록", "CU-04", "모달", "·"),
            ("엑셀(CSV) 고객 등록", "CU-05", "모달", "·"),
        ]]),
        ("고객불만", "CP-01", [[
            ("고객불만 상세", "CP-02", "드로어", None),
            ("고객불만 등록", "CP-03", "모달", "·"),
        ]]),
    ]),
    ("영업", None, [
        ("캘린더", "CA-01", [[
            ("일정 상세", "CA-02", "드로어", None),
            ("일정 등록", "CA-03", "모달", "·"),
        ]]),
        ("업무보고", "RP-01", [
            [
                ("보고서 대상 일정 고르기", "RP-07", None, None),
                ("업무보고서 작성", "RP-02", None, "→"),
                ("업무보고서 AI 결과", "RP-03", None, "→"),
                ("업무보고서 상세", "RP-08", None, "→"),
            ],
            [
                ("일일업무보고 작성", "RP-04", None, None),
                ("주간업무보고 작성", "RP-05", None, "→"),
                ("월간업무보고 작성", "RP-06", None, "→"),
            ],
        ]),
        ("영업현황", "DL-01 · DL-02", [[
            ("영업 딜 추가", "DL-03", "모달", None),
            ("딜 상세", "DL-04", "드로어", "·"),
            ("견적 현황", "QT-01", "시나리오 외", "→"),
            ("계약 현황", "CT-01", "시나리오 외", "·"),
            ("발주 관리", "OD-01", "시나리오 외", "·"),
        ]]),
        ("매출분석", "SA-01", []),
    ]),
    ("자료실", None, [
        ("자료실", "DC-01", [[
            ("파일 업로드", "DC-02", "모달", None),
            ("자료 상세", "DC-04", "드로어", "·"),
        ]]),
    ]),
    ("관리", "팀장 전용", [
        ("__MANAGER__", None, []),
    ]),
]

MANAGER = [("공지관리", "MG-01"), ("상품관리", "MG-02"), ("팀 관리", "MG-03")]

BAND = [
    ("인증 (앱 셸 밖)",
     "AU-01 로그인  →  AU-02 회원가입 · AU-03 비밀번호 설정"),
    ("사이드바 밖 진입",
     "NT-01 알림 (헤더 벨) · MY-01 마이페이지 (사이드바 하단 이름 · 헤더 아바타)"),
    ("공통 오류",
     "ER-01 404 · ER-02 서버 연결 불가 (전 화면 공통)"),
]


def card(slide, x, y, w, h, label, sid, tag=None,
         fill=WHITE, fg=INK, border=LINE, size=8):
    """Flow 슬라이드의 node() 와 같은 표기 규칙 — 이름 + 그 아래 화면 ID."""
    rect(slide, x, y, w, h, fill=fill, line=border, lw=1,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    sub = sid or ""
    if tag:
        sub = f"{sub} · {tag}" if sub else tag
    blocks = [(label, size, True, fg)]
    if sub:
        blocks.append((sub, 6.5, False, MUTED if fg == INK else SKY, 1.0, 1.5))
    text(slide, x, y, w, h, blocks, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def wire(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = WIRE
    c.line.width = Pt(1.25)
    return c


def glyph(slide, x, y, mark):
    text(slide, x, y, GAP, CARD_H, [(mark, 7, True, WIRE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def draw(slide):
    section_head(slide, "04", "사이트맵")
    text(slide, Inches(0.45), Inches(1.00), Inches(12.4), Inches(0.24),
         [("사이드바 메뉴를 기준으로 한 화면 계층입니다. "
           "화면 목록(02)의 45개 화면을 대분류 → 메뉴 → 그 메뉴에서 열리는 화면 순으로 놓았습니다. "
           "→ 는 이어지는 흐름, · 는 나란한 화면입니다.", 9, False, BODY)])

    y = Y0
    for title, sub, rows in SECTIONS:
        top = y
        for menu, mid, lines in rows:
            h = ROW if len(lines) <= 1 else ROW * len(lines) + VGAP * (len(lines) - 1)

            if menu == "__MANAGER__":
                # 관리 메뉴 셋은 하위 화면이 없어 한 줄에 나란히 폅니다.
                wire(slide, T1_X + T1_W, Inches(y + ROW / 2), T2_X, Inches(y + ROW / 2))
                for i, (label, sid) in enumerate(MANAGER):
                    cx = T2_X if i == 0 else T3_X + (CARD_W + GAP) * (i - 1)
                    if i:
                        glyph(slide, cx - GAP, Inches(y), "·")
                    card(slide, cx, Inches(y), T2_W if i == 0 else CARD_W, CARD_H,
                         label, sid, fill=TINT, size=9)
                y += h + VGAP
                continue

            card(slide, T2_X, Inches(y), T2_W, Inches(h), menu, mid,
                 fill=TINT, size=9)
            wire(slide, T1_X + T1_W, Inches(y + h / 2), T2_X, Inches(y + h / 2))

            for li, cards in enumerate(lines):
                ly = y + (ROW + VGAP) * li
                x = T3_X
                for ci, (label, sid, tag, sep) in enumerate(cards):
                    if ci == 0:
                        wire(slide, T2_X + T2_W, Inches(ly + ROW / 2),
                             x, Inches(ly + ROW / 2))
                    else:
                        glyph(slide, x - GAP, Inches(ly), sep or "·")
                    card(slide, x, Inches(ly), CARD_W, CARD_H, label, sid, tag)
                    x += CARD_W + GAP

            if not lines:
                text(slide, T3_X, Inches(y), T3_W, CARD_H,
                     [("이 메뉴에서 바로 열리는 하위 화면이 없습니다.", 8, False, MUTED)],
                     anchor=MSO_ANCHOR.MIDDLE)

            y += h + VGAP

        # 대분류는 그 아래 행 전체를 덮습니다.
        card(slide, T1_X, Inches(top), T1_W, Inches(y - VGAP - top), title, None, sub,
             fill=BRAND_DK, fg=WHITE, border=BRAND_DK, size=9.5)

    # 셸 밖 화면은 계층에 걸리지 않아 아래 띠로 뺍니다.
    by = Inches(6.28)
    rect(slide, T1_X, by, Inches(12.43), Inches(0.70), fill=TINT,
         line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    bw = Inches(4.00)
    for i, (cap, body) in enumerate(BAND):
        bx = T1_X + Inches(0.18) + (bw + Inches(0.10)) * i
        text(slide, bx, by + Inches(0.11), bw, Inches(0.20), [(cap, 8, True, BRAND)])
        text(slide, bx, by + Inches(0.35), bw, Inches(0.24), [(body, 7.5, False, BODY)])
        if i:
            rect(slide, bx - Inches(0.10), by + Inches(0.14),
                 Inches(0.010), Inches(0.42), fill=LINE)

    footer(slide, "사이트맵", PAGE)


def main():
    prs = Presentation(str(SRC))
    before = len(prs.slides._sldIdLst)
    draw(prs.slides.add_slide(prs.slide_layouts[6]))

    lst = prs.slides._sldIdLst
    added = list(lst)[-1]
    lst.remove(added)
    lst.insert(AT, added)

    prs.save(str(OUT))
    print(f"저장: {OUT}  ({before} → {len(lst)}장, {AT + 1}번째에 삽입)")


if __name__ == "__main__":
    main()
